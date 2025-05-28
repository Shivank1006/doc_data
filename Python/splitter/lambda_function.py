import json
import boto3
import aioboto3  # New import for async S3 operations
import asyncio   # New import for async operations
from pathlib import Path
import os
import logging
import uuid
import subprocess
from typing import Dict, List, Optional, Tuple
import shutil
from botocore.exceptions import BotoCoreError, ClientError  # New import for error handling

# Document processing libraries - ensure these are packaged properly in Lambda
import fitz  # PyMuPDF for PDF text extraction
from pdf2image import convert_from_path, pdfinfo_from_path
from pptx import Presentation
import docx2txt
from PIL import Image

# Configure logging
logger = logging.getLogger()
logger.setLevel(logging.INFO)

# --- S3Utils Class (copied from s3_utils.py and adapted) ---
class S3Utils:
    def __init__(self, bucket_name, access_key, secret_key, region_name, session_token=None) -> None:
        self.bucket_name = bucket_name
        self.access_key = access_key
        self.secret_key = secret_key
        self.region_name = region_name
        self.session_token = session_token
        # Create an aioboto3 session with credentials.
        self.session = aioboto3.Session(
            aws_access_key_id=self.access_key,
            aws_secret_access_key=self.secret_key,
            aws_session_token=self.session_token,
            region_name=self.region_name
        )
        logger.info(f"S3Utils initialized with bucket: {bucket_name}, region: {region_name}")

    async def read_file_from_s3_path(self, s3_path, local_filename):
        """
        Downloads a file from the specified S3 path to a local file path.
        Returns the local filename on success.
        """
        logger.info(f"Start downloading file from S3: {s3_path} to local path: {local_filename}")
        try:
            async with self.session.client('s3') as s3_client:
                await s3_client.download_file(self.bucket_name, s3_path, local_filename)
            logger.info(f"Finished downloading file from S3: {s3_path} to local path: {local_filename}")
            return local_filename
        except Exception as e:
            logger.error(f"Error downloading file from S3: {s3_path} to {local_filename}: {e}")
            raise

    async def write_file_to_s3(self, local_file_path, s3_file_path_with_extension):
        """
        Uploads a local file to the specified S3 path.
        Returns the S3 path on success.
        """
        logger.info(f"Start uploading local file {local_file_path} to S3: {s3_file_path_with_extension}")
        try:
            async with self.session.client('s3') as s3_client:
                await s3_client.upload_file(local_file_path, self.bucket_name, s3_file_path_with_extension)
            logger.info(f"Finished uploading local file {local_file_path} to S3: {s3_file_path_with_extension}")
            return f"s3://{self.bucket_name}/{s3_file_path_with_extension}"
        except Exception as e:
            logger.error(f"Error uploading file to S3: {local_file_path} to {s3_file_path_with_extension}: {e}")
            raise

# --- Read AWS credentials and environment variables ---
aws_access_key_id_env = os.environ.get('AWS_ACCESS_KEY_ID')
aws_secret_access_key_env = os.environ.get('AWS_SECRET_ACCESS_KEY')
aws_session_token_env = os.environ.get('AWS_SESSION_TOKEN') # This might be None
aws_region_env = os.environ.get('AWS_REGION')

# We'll keep the boto3 client for operations that don't use S3Utils
s3_client_config = {}

if aws_region_env:
    s3_client_config['region_name'] = aws_region_env

if aws_access_key_id_env and aws_secret_access_key_env:
    s3_client_config['aws_access_key_id'] = aws_access_key_id_env
    s3_client_config['aws_secret_access_key'] = aws_secret_access_key_env
    if aws_session_token_env:
        s3_client_config['aws_session_token'] = aws_session_token_env

s3_client = boto3.client('s3', **s3_client_config)
logger.info(f"Boto3 S3 client initialized. Region: {s3_client_config.get('region_name')}, Explicit creds used: {bool(s3_client_config.get('aws_access_key_id'))}")

try:
    BUCKET_NAME = os.environ['S3_BUCKET_NAME']
    INTERMEDIATE_IMAGES_PREFIX = os.environ.get('INTERMEDIATE_IMAGES_PREFIX', 'intermediate-images')
    INTERMEDIATE_RAW_TEXT_PREFIX = os.environ.get('INTERMEDIATE_RAW_TEXT_PREFIX', 'intermediate-raw-text')
    # Optional configuration for PDF to image conversion
    PDF_DPI = int(os.environ.get('PDF_DPI', '200'))
except KeyError as e:
    logger.error(f"Missing necessary environment variable: {e}")
    BUCKET_NAME = None # Indicate failure

# --- Helper Functions ---

def get_doc_type(file_extension: str) -> str:
    """Determine document type from file extension."""
    ext = file_extension.lower()
    if ext == ".pdf":
        return "pdf"
    elif ext in [".docx", ".doc"]:
        return "docx"
    elif ext in [".pptx", ".ppt"]:
        return "pptx"
    elif ext in [".png", ".jpg", ".jpeg"]:
        return "image"
    else:
        return "unsupported"

async def async_upload_files_to_s3(s3_utils, local_paths: List[str], s3_prefix: str) -> List[str]:
    """
    Uploads files from local paths to S3 under the specified prefix using S3Utils.
    This is an async version of the original upload_files_to_s3 function.
    """
    s3_uris = []
    if not local_paths:
        return s3_uris
    
    for local_path_str in local_paths:
        local_path = Path(local_path_str)
        if not local_path.is_file():
            logger.warning(f"Local file not found for upload: {local_path}, skipping.")
            continue
            
        s3_key = f"{s3_prefix}{local_path.name}" # Prefix should end with /
        try:
            logger.info(f"Uploading {local_path} to s3://{s3_utils.bucket_name}/{s3_key}")
            s3_uri = await s3_utils.write_file_to_s3(str(local_path), s3_key)
            s3_uris.append(s3_uri)
        except Exception as e:
            logger.error(f"Error uploading {local_path} to {s3_key}: {e}")
            # Continue with other files, just log the error
    
    return s3_uris

# The original synchronous upload_files_to_s3 function is kept as a wrapper around the async version
def upload_files_to_s3(local_paths: List[str], bucket: str, s3_prefix: str) -> List[str]:
    """Uploads files from local paths to S3 under the specified prefix using the async S3Utils."""
    # Create an S3Utils instance for this operation
    s3_utils = S3Utils(
        bucket_name=bucket,
        access_key=aws_access_key_id_env,
        secret_key=aws_secret_access_key_env,
        region_name=aws_region_env,
        session_token=aws_session_token_env
    )
    
    # Run the async function using asyncio.run
    return asyncio.run(async_upload_files_to_s3(s3_utils, local_paths, s3_prefix))

def _run_soffice(cmd: List[str]) -> None:
    """Run LibreOffice command with error handling."""
    try:
        subprocess.run(cmd, check=True)
    except subprocess.CalledProcessError as e:
        logger.error(f"LibreOffice conversion error: {e}")
        raise
    except FileNotFoundError:
        logger.error("Error: LibreOffice (soffice) not found. Please ensure LibreOffice is installed and in PATH.")
        raise

def save_content_to_local_file(content: str, output_path: Path) -> bool:
    """Saves string content to a local file in /tmp."""
    try:
        output_path.parent.mkdir(parents=True, exist_ok=True)
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(content)
        logger.info(f"Saved file: {output_path}")
        return True
    except IOError as e:
        logger.error(f"Error saving file {output_path}: {e}")
        return False
    except Exception as e:
        logger.error(f"Unexpected error saving file {output_path}: {e}")
        return False

# === PDF Handling Functions ===

def extract_text_from_pdf(pdf_path: str, text_output_dir: Path, base_filename: str) -> List[str]:
    """
    Extracts text from PDF pages using PyMuPDF (fitz), saving each page's text to a file.
    Returns list of paths to page text files.
    """
    logger.info(f"Extracting text from PDF: {pdf_path}")
    text_output_dir.mkdir(parents=True, exist_ok=True)
    page_text_paths = []

    try:
        doc = fitz.open(pdf_path)
        for page_num in range(doc.page_count):
            page = doc.load_page(page_num)
            page_text = page.get_text("text")
            
            if page_text.strip():  # Only save non-empty text
                text_filename = f"{base_filename}_page_{page_num+1}_text.txt"
                text_path = text_output_dir / text_filename
                
                if save_content_to_local_file(page_text, text_path):
                    page_text_paths.append(str(text_path))
                    logger.info(f"Saved text for page {page_num+1} to {text_path}")
                else:
                    logger.warning(f"Failed to save text for page {page_num+1}")
        
        logger.info(f"Text extraction from PDF complete. Extracted {len(page_text_paths)} page text files.")
        return page_text_paths
    except Exception as e:
        logger.error(f"Error extracting text from PDF {pdf_path}: {e}")
        return []

def convert_pdf_to_images(pdf_path: str, image_output_dir: Path, base_filename: str, dpi: int = PDF_DPI) -> List[str]:
    """
    Converts PDF pages to images using pdf2image.
    Returns list of paths to page images.
    """
    logger.info(f"Converting PDF to images: {pdf_path} with DPI={dpi}")
    image_output_dir.mkdir(parents=True, exist_ok=True)
    image_paths = []
    
    try:
        # Get PDF info for logging
        try:
            info = pdfinfo_from_path(pdf_path)
            page_count = info.get('Pages', 'unknown')
            logger.info(f"PDF Info: {page_count} pages detected.")
        except Exception as info_err:
            logger.warning(f"Could not get PDF info: {info_err}. Proceeding.")
        
        # Convert PDF to images
        images = convert_from_path(pdf_path, dpi=dpi)
        if not images:
            logger.warning(f"Warning: pdf2image returned no images for {pdf_path}.")
            return []
        
        logger.info(f"Conversion complete. Found {len(images)} images. Saving locally...")
        
        for i, img in enumerate(images):
            page_num = i + 1
            img_filename = f"{base_filename}_page_{page_num}.png"
            local_img_path = image_output_dir / img_filename
            
            try:
                img.save(local_img_path, "PNG")
                logger.info(f"Saved image: {local_img_path}")
                image_paths.append(str(local_img_path))
            except Exception as page_save_err:
                logger.error(f"Error saving page {page_num} image to {local_img_path}: {page_save_err}")
        
        logger.info(f"PDF to image conversion finished. Saved {len(image_paths)} images.")
        return image_paths
    except Exception as e:
        logger.error(f"Error converting PDF to images: {e}")
        return []

# === DOCX Handling Functions ===

def extract_text_from_docx(docx_path: str, text_output_dir: Path, base_filename: str) -> List[str]:
    """
    Extracts full text from DOCX using docx2txt.
    Returns a list with a single path to the text file.
    """
    logger.info(f"Extracting text from DOCX: {docx_path}")
    text_output_dir.mkdir(parents=True, exist_ok=True)
    
    try:
        # Extract all text using docx2txt
        docx_text = docx2txt.process(docx_path)
        
        if docx_text.strip():
            text_filename = f"{base_filename}_full_text.txt"
            text_path = text_output_dir / text_filename
            
            if save_content_to_local_file(docx_text, text_path):
                logger.info(f"Saved DOCX text to {text_path}")
                return [str(text_path)]
            else:
                logger.warning("Failed to save DOCX text")
        else:
            logger.warning("No text content found in DOCX")
        
        return []
    except Exception as e:
        logger.error(f"Error extracting text from DOCX {docx_path}: {e}")
        return []

def convert_office_to_pdf(office_path: str, temp_dir: str) -> str:
    """
    Converts Office document (DOCX/PPTX) to PDF using LibreOffice.
    Returns path to the generated PDF.
    """
    input_path = Path(office_path)
    base_name = input_path.stem
    output_pdf = os.path.join(temp_dir, f"{base_name}.pdf")
    
    logger.info(f"Converting {office_path} to PDF using LibreOffice...")
    cmd = ["soffice", "--headless", "--convert-to", "pdf", 
           "--outdir", temp_dir, office_path]
    
    _run_soffice(cmd)
    
    if os.path.exists(output_pdf):
        logger.info(f"Successfully converted to PDF: {output_pdf}")
        return output_pdf
    else:
        raise FileNotFoundError(f"LibreOffice conversion failed. PDF not found: {output_pdf}")

def convert_office_to_images(office_path: str, image_output_dir: Path, base_filename: str, temp_dir: str) -> List[str]:
    """
    Converts Office document (DOCX/PPTX) to images via PDF using LibreOffice.
    Returns list of paths to page images.
    """
    image_output_dir.mkdir(parents=True, exist_ok=True)
    
    try:
        # Step 1: Convert Office document to PDF
        pdf_path = convert_office_to_pdf(office_path, temp_dir)
        
        # Step 2: Convert PDF to images
        return convert_pdf_to_images(pdf_path, image_output_dir, base_filename)
    except Exception as e:
        logger.error(f"Error converting office document to images: {e}")
        return []

# === PPTX Handling Functions ===

def extract_text_from_pptx(pptx_path: str, text_output_dir: Path, base_filename: str) -> List[str]:
    """
    Extracts text from PPTX slides, saving each slide's text to a file.
    Returns list of paths to slide text files.
    """
    logger.info(f"Extracting text from PPTX: {pptx_path}")
    text_output_dir.mkdir(parents=True, exist_ok=True)
    page_text_paths = []
    
    try:
        prs = Presentation(pptx_path)
        for slide_num, slide in enumerate(prs.slides):
            # Collect text from all shapes in the slide
            text_items = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_items.append(shape.text.strip())
            
            slide_text = "\n".join(text_items)
            if slide_text:  # Only save non-empty text
                text_filename = f"{base_filename}_slide_{slide_num+1}_text.txt"
                text_path = text_output_dir / text_filename
                
                if save_content_to_local_file(slide_text, text_path):
                    page_text_paths.append(str(text_path))
                    logger.info(f"Saved text for slide {slide_num+1} to {text_path}")
                else:
                    logger.warning(f"Failed to save text for slide {slide_num+1}")
        
        logger.info(f"Text extraction from PPTX complete. Extracted {len(page_text_paths)} slide text files.")
        return page_text_paths
    except Exception as e:
        logger.error(f"Error extracting text from PPTX {pptx_path}: {e}")
        return []

# === Image Handling Functions ===

def process_input_image(image_path: str, image_output_dir: Path, base_filename: str) -> List[str]:
    """
    Processes image input file: standardizes to PNG, saves to output directory.
    Returns list containing path to the processed image.
    """
    logger.info(f"Processing image input: {image_path}")
    image_output_dir.mkdir(parents=True, exist_ok=True)
    
    try:
        img = Image.open(image_path)
        img_filename = f"{base_filename}_image.png"
        output_img_path = image_output_dir / img_filename
        
        img.save(output_img_path, "PNG")
        logger.info(f"Saved standardized image: {output_img_path}")
        return [str(output_img_path)]
    except Exception as e:
        logger.error(f"Error processing input image {image_path}: {e}")
        return []

# === Main Splitter Logic Function ===

def run_splitter_logic(
    source_file_path: str,
    file_extension: str,
    base_filename: str,
    run_uuid: str,
    local_temp_dir: str,
    local_image_output_dir: Path,
    local_text_output_dir: Path
) -> Tuple[List[str], List[str]]:
    """
    Adapted core splitter logic for Lambda.
    It performs text extraction and image conversion,
    saving results to the specified local /tmp directories.

    Args:
        source_file_path: Path to the downloaded input file in /tmp.
        file_extension: The extension of the input file (e.g., '.pdf').
        base_filename: Original base filename (used for naming local output files).
        run_uuid: The unique ID for this run.
        local_temp_dir: Path to the main temporary directory for this run in /tmp.
        local_image_output_dir: Path object for saving generated images in /tmp.
        local_text_output_dir: Path object for saving extracted text in /tmp.

    Returns:
        A tuple containing: (list of local text file paths, list of local image file paths)
    """
    logger.info(f"Running splitter logic for: {source_file_path} (type based on extension: {file_extension})")
    
    doc_type = get_doc_type(file_extension)
    page_text_paths = []
    page_image_paths = []

    # Process based on document type
    try:
        # 1. Extract text (if applicable)
        if doc_type == "pdf":
            page_text_paths = extract_text_from_pdf(source_file_path, local_text_output_dir, base_filename)
        elif doc_type == "docx":
            page_text_paths = extract_text_from_docx(source_file_path, local_text_output_dir, base_filename)
        elif doc_type == "pptx":
            page_text_paths = extract_text_from_pptx(source_file_path, local_text_output_dir, base_filename)
        # No text extraction for images

        # 2. Convert to images
        if doc_type == "pdf":
            page_image_paths = convert_pdf_to_images(source_file_path, local_image_output_dir, base_filename)
        elif doc_type in ["docx", "pptx"]:
            page_image_paths = convert_office_to_images(source_file_path, local_image_output_dir, base_filename, local_temp_dir)
        elif doc_type == "image":
            page_image_paths = process_input_image(source_file_path, local_image_output_dir, base_filename)

    except Exception as e:
        logger.error(f"Error processing {source_file_path}: {e}")
        import traceback
        logger.error(traceback.format_exc())
        return [], []

    if not page_image_paths:
        logger.warning(f"Warning: No images were generated for {source_file_path}.")
    
    if not page_text_paths and doc_type != "image":
        logger.warning(f"Warning: No text was extracted for {source_file_path}.")

    logger.info(f"Splitter logic finished. Generated {len(page_image_paths)} images and {len(page_text_paths)} text files.")
    return page_text_paths, page_image_paths

# --- Main Lambda Handler ---

def lambda_handler(event, context):
    """
    AWS Lambda handler for the Splitter function.
    Downloads the input file, splits it into images and text (per page),
    and uploads the results back to S3.
    Input event (example):
    {'s3_input_uri': 's3://bucket/inputs/mydoc.pdf', 'output_format': 'markdown'}
    Output: JSON containing S3 URIs for page images and text files.
    """
    logger.info(f"Received event: {json.dumps(event)}")

    if not BUCKET_NAME:
        # Logged the error during initialization
        return {'statusCode': 500, 'body': json.dumps('Error: S3_BUCKET_NAME not configured.')}

    try:
        # s3_input_uri from the event is now expected to be the S3 object key (e.g., "path/to/file.pdf")
        input_key = event['s3_input_uri']
        output_format = event['output_format'] # Keep track of the desired final format
        if not input_key or not isinstance(input_key, str):
            raise KeyError("s3_input_uri (object key) must be a non-empty string.")
    except KeyError as e:
        logger.error(f"Missing key in input event or invalid key: {e}")
        return {'statusCode': 400, 'body': json.dumps(f'Error: Missing or invalid required key: {e}')}

    # --- 1. Parse S3 Key and Generate UUID ---
    try:
        if input_key.startswith('s3://'):
            logger.error(f"Invalid s3_input_uri format. Expected S3 object key, got full URI: {input_key}")
            raise ValueError("s3_input_uri should be an S3 object key (e.g., 'path/to/file.pdf'), not a full S3 URI.")

        original_filename = Path(input_key).name
        base_filename = Path(original_filename).stem # Keep original base for reference/naming
        file_extension = Path(original_filename).suffix.lower()
        run_uuid = str(uuid.uuid4()) # Generate UUID for this run
        logger.info(f"Starting splitter run {run_uuid} for {original_filename} (S3 Key: {input_key})")

        doc_type = get_doc_type(file_extension)
        if doc_type == "unsupported":
            raise ValueError(f"Unsupported file type: {file_extension}")

    except ValueError as e:
        logger.error(f"Invalid input: {e}")
        return {'statusCode': 400, 'body': json.dumps(f'Invalid input: {e}')}
    except Exception as e:
        logger.error(f"Error during input parsing: {e}")
        return {'statusCode': 500, 'body': json.dumps(f'Error parsing input: {e}')}


    # --- 2. Download Input File to /tmp ---
    local_input_path = Path('/tmp') / f"{run_uuid}_{original_filename}" # Add UUID to avoid collisions
    try:
        # Create an S3Utils instance for this operation
        s3_utils = S3Utils(
            bucket_name=BUCKET_NAME,
            access_key=aws_access_key_id_env,
            secret_key=aws_secret_access_key_env,
            region_name=aws_region_env,
            session_token=aws_session_token_env
        )
        
        logger.info(f"Downloading {input_key} to {local_input_path} using S3Utils...")
        # Use asyncio.run to call the async method
        asyncio.run(s3_utils.read_file_from_s3_path(input_key, str(local_input_path)))
        logger.info("Download complete.")
    except Exception as e:
        logger.error(f"Error downloading {input_key}: {e}")
        return {'statusCode': 500, 'body': json.dumps(f'Error downloading input file: {e}')}

    # --- 3. Setup Local Directories & Run Core Logic ---
    # Use UUID for run-specific temp directory
    local_temp_dir = Path('/tmp') / run_uuid
    local_image_output_dir = local_temp_dir / "images"
    local_text_output_dir = local_temp_dir / "text"
    try:
        local_image_output_dir.mkdir(parents=True, exist_ok=True)
        local_text_output_dir.mkdir(parents=True, exist_ok=True)
        logger.info(f"Created temporary directories under {local_temp_dir}")

        # Call adapted core splitter logic
        page_text_local_paths, page_image_local_paths = run_splitter_logic(
             source_file_path=str(local_input_path),
             file_extension=file_extension,
             base_filename=base_filename, # Still useful for local file naming
             run_uuid=run_uuid,
             local_temp_dir=str(local_temp_dir),
             local_image_output_dir=local_image_output_dir,
             local_text_output_dir=local_text_output_dir
        )
        logger.info(f"Core splitter logic finished. Found {len(page_text_local_paths)} text files, {len(page_image_local_paths)} images.")

    except Exception as e:
        logger.error(f"Error during core splitting process: {e}")
        # Perform cleanup before returning error
        # import shutil
        # shutil.rmtree(local_temp_dir, ignore_errors=True)
        # local_input_path.unlink(missing_ok=True)
        return {'statusCode': 500, 'body': json.dumps(f'Error during splitting: {e}')}

    # --- 4. Upload Results to S3 ---
    # Define S3 output prefixes using UUID. Ensure they end with /
    s3_image_output_prefix = f"{INTERMEDIATE_IMAGES_PREFIX}/{run_uuid}/"
    s3_text_output_prefix = f"{INTERMEDIATE_RAW_TEXT_PREFIX}/{run_uuid}/"

    s3_page_text_uris = upload_files_to_s3(page_text_local_paths, BUCKET_NAME, s3_text_output_prefix)
    s3_page_image_uris = upload_files_to_s3(page_image_local_paths, BUCKET_NAME, s3_image_output_prefix)
    logger.info(f"Uploaded {len(s3_page_text_uris)} text URIs, {len(s3_page_image_uris)} image URIs.")

    # --- 5. Clean up /tmp (Best effort) ---
    try:
        # Remove specific files first
        local_input_path.unlink(missing_ok=True)
        for p in page_text_local_paths: Path(p).unlink(missing_ok=True)
        for p in page_image_local_paths: Path(p).unlink(missing_ok=True)
        # Remove the run-specific directory
        if local_temp_dir.exists():
            import shutil
            shutil.rmtree(local_temp_dir, ignore_errors=True)
            logger.info(f"Cleaned up temporary directory: {local_temp_dir}")
    except Exception as e:
        logger.warning(f"Error during /tmp cleanup: {e}")

    # --- 6. Return S3 URIs and Metadata for Step Functions ---
    # This output becomes the input for the Map state iterating over pages
    reconstructed_s3_uri = f"s3://{BUCKET_NAME}/{input_key}"
    if doc_type == "image":
        s3_page_text_uris = [None] * len(s3_page_image_uris)
    result = {
        'run_uuid': run_uuid,
        'original_s3_uri': reconstructed_s3_uri, # Reconstructed full S3 URI
        'original_s3_key': input_key,           # Explicitly include the key used
        'original_base_filename': base_filename,
        'doc_type': doc_type,
        'output_format': output_format,
        's3_page_text_uris': s3_page_text_uris, # List of S3 URIs
        's3_page_image_uris': s3_page_image_uris, # List of S3 URIs
    }
    logger.info(f"Splitter finished successfully for run {run_uuid}. Output: {json.dumps(result)}")
    return result

# Example Test Event:
# {
#   "s3_input_uri": "inputs/example.pdf", <--- Note: This is now an S3 object key
#   "output_format": "markdown"
# } 